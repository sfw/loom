"""Image and document processing utilities for multimodal content.

Provides base64 encoding, image dimension parsing, resize support,
and PDF-to-image rendering. All heavy dependencies (Pillow, pymupdf)
are optional — functions degrade gracefully when unavailable.
"""

from __future__ import annotations

import base64
import logging
import struct
from pathlib import Path

logger = logging.getLogger(__name__)

MAX_IMAGE_BYTES = 5 * 1024 * 1024  # 5MB per image
MAX_PDF_BYTES = 32 * 1024 * 1024  # 32MB per PDF
IMAGE_MAX_DIMENSION = 2048


def encode_image_base64(path: Path, max_dimension: int = IMAGE_MAX_DIMENSION) -> str | None:
    """Read and base64-encode an image file.

    If the image exceeds max_dimension on either axis and Pillow is
    available, resize it proportionally. Without Pillow, returns the
    raw image if within the 5MB limit, or None if too large.
    """
    try:
        if not path.exists():
            return None

        size = path.stat().st_size
        if size == 0:
            return None
        if size > MAX_IMAGE_BYTES:
            try:
                return _resize_and_encode(path, max_dimension)
            except ImportError:
                return None

        data = path.read_bytes()

        w, h = get_image_dimensions_from_bytes(data[:64], path.suffix)
        if (w > max_dimension or h > max_dimension) and w > 0:
            try:
                return _resize_and_encode(path, max_dimension)
            except ImportError:
                pass  # Send as-is, let the API handle it

        return base64.b64encode(data).decode("ascii")
    except (OSError, PermissionError) as e:
        logger.debug("Failed to encode image %s: %s", path, e)
        return None


def encode_file_base64(path: Path) -> str | None:
    """Read and base64-encode any file (used for PDFs)."""
    try:
        if not path.exists():
            return None
        data = path.read_bytes()
        if not data:
            return None
        return base64.b64encode(data).decode("ascii")
    except (OSError, PermissionError) as e:
        logger.debug("Failed to encode file %s: %s", path, e)
        return None


def get_image_dimensions(path: Path) -> tuple[int, int]:
    """Read image dimensions from file header without loading the full image.

    Supports PNG, JPEG, GIF, BMP, WebP headers.
    Returns (0, 0) for unrecognized formats.
    """
    try:
        header = path.read_bytes()[:32]
        return get_image_dimensions_from_bytes(header, path.suffix)
    except (OSError, PermissionError):
        return (0, 0)


def get_image_dimensions_from_bytes(
    header: bytes, suffix: str,
) -> tuple[int, int]:
    """Parse image dimensions from header bytes."""
    try:
        # PNG
        if header[:8] == b"\x89PNG\r\n\x1a\n" and len(header) >= 24:
            w, h = struct.unpack(">II", header[16:24])
            return (w, h)
        # JPEG — complex header; skip for now
        if header[:2] == b"\xff\xd8":
            return (0, 0)
        # GIF
        if header[:6] in (b"GIF87a", b"GIF89a") and len(header) >= 10:
            w, h = struct.unpack("<HH", header[6:10])
            return (w, h)
        # BMP
        if header[:2] == b"BM" and len(header) >= 26:
            w, h = struct.unpack("<II", header[18:26])
            return (w, h)
        # WebP
        if header[:4] == b"RIFF" and len(header) >= 16 and header[8:12] == b"WEBP":
            if header[12:16] == b"VP8 " and len(header) >= 30:
                w = struct.unpack("<H", header[26:28])[0] & 0x3FFF
                h = struct.unpack("<H", header[28:30])[0] & 0x3FFF
                return (w, h)
    except struct.error:
        pass
    return (0, 0)


def _resize_and_encode(path: Path, max_dim: int) -> str:
    """Resize with Pillow and return base64. Raises ImportError if no Pillow."""
    import io

    from PIL import Image

    with Image.open(path) as img:
        img.thumbnail((max_dim, max_dim), Image.LANCZOS)
        buf = io.BytesIO()
        fmt = "PNG" if path.suffix.lower() == ".png" else "JPEG"
        img.save(buf, format=fmt, quality=85)
        return base64.b64encode(buf.getvalue()).decode("ascii")


def extract_docx_text(path: Path) -> str:
    """Extract text content from a .docx file.

    Extracts paragraph text in document order.
    Returns the extracted text or raises on failure.
    """
    from docx import Document

    doc = Document(str(path))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    return "\n\n".join(paragraphs)


def extract_pptx_text(path: Path) -> str:
    """Extract text content from a .pptx file.

    Iterates over slides and extracts text from all shapes,
    prefixed with slide numbers.
    Returns the extracted text or raises on failure.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides_text.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def pdf_pages_to_images(
    path: Path,
    page_range: tuple[int, int] | None = None,
    dpi: int = 150,
    max_dimension: int = IMAGE_MAX_DIMENSION,
) -> list[str]:
    """Render PDF pages as base64-encoded images.

    Used for Ollama (no native PDF support) with vision-capable models.
    Requires pymupdf or pdf2image. Returns empty list if neither is available.
    """
    try:
        import fitz

        doc = fitz.open(path)
        try:
            start, end = page_range or (0, len(doc))
            images = []
            for i in range(start, min(end, len(doc))):
                page = doc[i]
                mat = fitz.Matrix(dpi / 72, dpi / 72)
                pix = page.get_pixmap(matrix=mat)
                if pix.width > max_dimension or pix.height > max_dimension:
                    scale = max_dimension / max(pix.width, pix.height)
                    mat = fitz.Matrix(scale * dpi / 72, scale * dpi / 72)
                    pix = page.get_pixmap(matrix=mat)
                images.append(base64.b64encode(pix.tobytes("png")).decode("ascii"))
            return images
        finally:
            doc.close()
    except ImportError:
        pass

    try:
        import io

        from pdf2image import convert_from_path

        start, end = page_range or (0, 999)
        pages = convert_from_path(
            path, dpi=dpi, first_page=start + 1, last_page=end,
        )
        images = []
        for page in pages:
            page.thumbnail((max_dimension, max_dimension))
            buf = io.BytesIO()
            page.save(buf, format="PNG")
            images.append(base64.b64encode(buf.getvalue()).decode("ascii"))
        return images
    except ImportError:
        return []
