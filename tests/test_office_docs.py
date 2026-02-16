"""Tests for Word (.docx) and PowerPoint (.pptx) document support."""

from __future__ import annotations

from pathlib import Path

import pytest

from loom.content import DocumentBlock
from loom.content_utils import extract_docx_text, extract_pptx_text
from loom.tools.file_ops import ReadFileTool
from loom.tools.registry import ToolContext


@pytest.fixture()
def workspace(tmp_path: Path) -> Path:
    return tmp_path


@pytest.fixture()
def ctx(workspace: Path) -> ToolContext:
    return ToolContext(workspace=workspace)


@pytest.fixture()
def read_tool() -> ReadFileTool:
    return ReadFileTool()


# ---------------------------------------------------------------------------
# Helper: create minimal .docx and .pptx files using the libraries
# ---------------------------------------------------------------------------


def _create_docx(path: Path, paragraphs: list[str]) -> Path:
    from docx import Document

    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    doc.save(str(path))
    return path


def _create_pptx(path: Path, slides: list[list[str]]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for texts in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        for i, text in enumerate(texts):
            txbox = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(6), Inches(1))
            txbox.text_frame.text = text
    prs.save(str(path))
    return path


# ---------------------------------------------------------------------------
# extract_docx_text
# ---------------------------------------------------------------------------


class TestExtractDocxText:
    def test_basic_extraction(self, tmp_path: Path):
        path = _create_docx(tmp_path / "test.docx", ["Hello world", "Second paragraph"])
        text = extract_docx_text(path)
        assert "Hello world" in text
        assert "Second paragraph" in text

    def test_empty_document(self, tmp_path: Path):
        path = _create_docx(tmp_path / "empty.docx", [])
        text = extract_docx_text(path)
        assert text == ""

    def test_whitespace_only_paragraphs_skipped(self, tmp_path: Path):
        path = _create_docx(tmp_path / "ws.docx", ["Real text", "   ", "More text"])
        text = extract_docx_text(path)
        assert "Real text" in text
        assert "More text" in text
        # Whitespace-only paragraph should not produce a section
        sections = [s for s in text.split("\n\n") if s.strip()]
        assert len(sections) == 2


# ---------------------------------------------------------------------------
# extract_pptx_text
# ---------------------------------------------------------------------------


class TestExtractPptxText:
    def test_basic_extraction(self, tmp_path: Path):
        path = _create_pptx(
            tmp_path / "test.pptx",
            [["Title Slide", "Subtitle"], ["Slide 2 content"]],
        )
        text = extract_pptx_text(path)
        assert "Slide 1" in text
        assert "Title Slide" in text
        assert "Slide 2" in text
        assert "Slide 2 content" in text

    def test_empty_presentation(self, tmp_path: Path):
        from pptx import Presentation

        prs = Presentation()
        path = tmp_path / "empty.pptx"
        prs.save(str(path))
        text = extract_pptx_text(path)
        assert text == ""

    def test_slide_numbering(self, tmp_path: Path):
        path = _create_pptx(
            tmp_path / "numbered.pptx",
            [["First"], ["Second"], ["Third"]],
        )
        text = extract_pptx_text(path)
        assert "--- Slide 1 ---" in text
        assert "--- Slide 2 ---" in text
        assert "--- Slide 3 ---" in text


# ---------------------------------------------------------------------------
# ReadFileTool — DOCX
# ---------------------------------------------------------------------------


class TestReadFileToolDocx:
    @pytest.mark.asyncio
    async def test_read_docx(self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool):
        _create_docx(workspace / "report.docx", ["Executive summary", "Details here"])
        result = await read_tool.execute({"path": "report.docx"}, ctx)
        assert result.success
        assert "Executive summary" in result.output
        assert "Details here" in result.output
        assert result.content_blocks is not None
        assert len(result.content_blocks) == 1
        block = result.content_blocks[0]
        assert isinstance(block, DocumentBlock)
        assert block.extracted_text
        assert block.size_bytes > 0

    @pytest.mark.asyncio
    async def test_read_docx_data(self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool):
        _create_docx(workspace / "data.docx", ["Content"])
        result = await read_tool.execute({"path": "data.docx"}, ctx)
        assert result.data["type"] == "docx"
        assert result.data["name"] == "data.docx"

    @pytest.mark.asyncio
    async def test_read_doc_extension(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        """The .doc extension routes through the docx handler (best-effort)."""
        # We create a .docx but name it .doc — python-docx can still read it
        _create_docx(workspace / "legacy.doc", ["Legacy content"])
        result = await read_tool.execute({"path": "legacy.doc"}, ctx)
        assert result.success
        assert "Legacy content" in result.output

    @pytest.mark.asyncio
    async def test_read_docx_empty(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        _create_docx(workspace / "blank.docx", [])
        result = await read_tool.execute({"path": "blank.docx"}, ctx)
        assert result.success
        assert "no extractable text" in result.output

    @pytest.mark.asyncio
    async def test_read_docx_corrupt(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        (workspace / "bad.docx").write_bytes(b"not a zip file")
        result = await read_tool.execute({"path": "bad.docx"}, ctx)
        assert not result.success
        assert "Error reading Word document" in result.error


# ---------------------------------------------------------------------------
# ReadFileTool — PPTX
# ---------------------------------------------------------------------------


class TestReadFileToolPptx:
    @pytest.mark.asyncio
    async def test_read_pptx(self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool):
        _create_pptx(
            workspace / "deck.pptx",
            [["Welcome", "Introduction"], ["Main Point"]],
        )
        result = await read_tool.execute({"path": "deck.pptx"}, ctx)
        assert result.success
        assert "Welcome" in result.output
        assert "Main Point" in result.output
        assert result.content_blocks is not None
        block = result.content_blocks[0]
        assert isinstance(block, DocumentBlock)
        assert "presentationml" in block.media_type

    @pytest.mark.asyncio
    async def test_read_pptx_data(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        _create_pptx(workspace / "info.pptx", [["Slide"]])
        result = await read_tool.execute({"path": "info.pptx"}, ctx)
        assert result.data["type"] == "pptx"
        assert result.data["name"] == "info.pptx"

    @pytest.mark.asyncio
    async def test_read_ppt_extension(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        """The .ppt extension routes through the pptx handler."""
        _create_pptx(workspace / "old.ppt", [["Old format"]])
        result = await read_tool.execute({"path": "old.ppt"}, ctx)
        assert result.success
        assert "Old format" in result.output

    @pytest.mark.asyncio
    async def test_read_pptx_empty(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        from pptx import Presentation

        prs = Presentation()
        prs.save(str(workspace / "empty.pptx"))
        result = await read_tool.execute({"path": "empty.pptx"}, ctx)
        assert result.success
        assert "no extractable text" in result.output

    @pytest.mark.asyncio
    async def test_read_pptx_corrupt(
        self, workspace: Path, ctx: ToolContext, read_tool: ReadFileTool
    ):
        (workspace / "bad.pptx").write_bytes(b"not a zip")
        result = await read_tool.execute({"path": "bad.pptx"}, ctx)
        assert not result.success
        assert "Error reading PowerPoint" in result.error
